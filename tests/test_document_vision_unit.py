from __future__ import annotations

from io import BytesIO
from types import SimpleNamespace

from marginalia.pipelines.document_vision import (
    DocumentImage,
    inline_document_image_vision_text,
    persisted_document_image_segment,
    select_document_images,
)
from marginalia.pipelines.docx import _extract_docx_images
from marginalia.pipelines.pptx import _extract_pptx_images


def _png_bytes() -> bytes:
    from PIL import Image

    image = Image.new("RGB", (96, 64), color=(220, 40, 40))
    out = BytesIO()
    image.save(out, format="PNG")
    return out.getvalue()


def test_select_document_images_scopes_by_slide_and_index() -> None:
    images = [
        DocumentImage(
            image_id="pptx-img-1",
            label="one",
            image_bytes=b"1",
            media_type="image/png",
            anchor={"slide": 1},
        ),
        DocumentImage(
            image_id="pptx-img-2",
            label="two",
            image_bytes=b"2",
            media_type="image/png",
            anchor={"slide": 3},
        ),
    ]

    assert [img.image_id for img in select_document_images(images, {"slide_start": 3})] == [
        "pptx-img-2"
    ]
    assert [img.image_id for img in select_document_images(images, {"image_index": 1})] == [
        "pptx-img-1"
    ]
    assert [img.image_id for img in select_document_images(images, {"image_id": "pptx-img-2"})] == [
        "pptx-img-2"
    ]


def test_persisted_document_image_segment_reads_description_payload() -> None:
    file_row = SimpleNamespace(
        description={
            "document_vision": {
                "text": "[DOCX image 1]\nVisible text: ACME contract",
            }
        }
    )

    result = persisted_document_image_segment(
        file_row,
        {"max_chars": 80},
        mode="docx_image_vision",
    )

    assert result.error is None
    assert "ACME" in result.text
    assert result.extras["source"] == "persisted_document_image_vision"


def test_inline_document_image_vision_text_anchors_to_slide() -> None:
    payload = {
        "images": [
            {
                "image_id": "pptx-img-1",
                "label": "PPTX slide 2 image 1",
                "anchor": {"slide": 2},
                "context": "Slide title",
                "text": "Visible text: Q3 revenue chart.",
            }
        ],
        "text": "[PPTX slide 2 image 1]\nVisible text: Q3 revenue chart.",
    }

    slides = inline_document_image_vision_text(
        ["# Slide 1\nIntro", "# Slide 2\nFinancials"],
        payload,
        anchor_key="slide",
    )

    assert "Q3 revenue chart" not in slides[0]
    assert "Embedded image vision" in slides[1]
    assert "Q3 revenue chart" in slides[1]


def test_inline_document_image_vision_text_anchors_to_docx_block() -> None:
    payload = {
        "images": [
            {
                "image_id": "docx-img-1",
                "label": "DOCX image 1",
                "anchor": {"block": 1},
                "text": "Visible text: signed approval stamp.",
            }
        ]
    }

    paragraphs = inline_document_image_vision_text(
        ["Approval memo", "Next section"],
        payload,
        anchor_key="block",
    )

    assert "signed approval stamp" in paragraphs[0]
    assert "signed approval stamp" not in paragraphs[1]


def test_extract_docx_images_from_embedded_picture() -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Context near the chart.")
    image_path = BytesIO(_png_bytes())
    doc.add_picture(image_path)
    out = BytesIO()
    doc.save(out)

    images = _extract_docx_images(out.getvalue())

    assert len(images) == 1
    assert images[0].image_id == "docx-img-1"
    assert images[0].media_type == "image/png"
    assert images[0].anchor["block"] == 1
    assert images[0].anchor["source_block"] == 2


def test_extract_pptx_images_from_embedded_picture(tmp_path) -> None:
    from pptx import Presentation

    image_file = tmp_path / "image.png"
    image_file.write_bytes(_png_bytes())
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_picture(str(image_file), 0, 0)
    out = BytesIO()
    deck.save(out)

    images = _extract_pptx_images(out.getvalue())

    assert len(images) == 1
    assert images[0].image_id == "pptx-img-1"
    assert images[0].anchor["slide"] == 1
